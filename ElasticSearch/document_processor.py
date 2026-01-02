"""
Document Processor for Elasticsearch Multi-Document Search System

Extracts text from PDF, DOCX, TXT, CSV, and Excel files
Uses PyMuPDF for better PDF extraction and RecursiveCharacterTextSplitter for semantic chunking

"""

import os
import sys
import importlib.util
from typing import List, Dict, Any
from pathlib import Path
import fitz
import uuid
import hashlib  # Using SHA256 for document hashing
from docx import Document
import pandas as pd
from pptx import Presentation
from langchain_text_splitters import RecursiveCharacterTextSplitter

# Add ElasticSearch directory to path and load local utils
_current_dir = os.path.dirname(os.path.abspath(__file__))
_project_root = os.path.dirname(_current_dir)
if _current_dir not in sys.path:
    sys.path.insert(0, _current_dir)
if _project_root not in sys.path:
    sys.path.insert(0, _project_root)

# Load local ElasticSearch/utils.py using importlib to avoid conflict with utils/ package
_es_utils_path = os.path.join(_current_dir, "utils.py")
_es_utils_spec = importlib.util.spec_from_file_location("es_utils", _es_utils_path)
_es_utils = importlib.util.module_from_spec(_es_utils_spec)
_es_utils_spec.loader.exec_module(_es_utils)

setup_logging = _es_utils.setup_logging
clean_text = _es_utils.clean_text
count_words = _es_utils.count_words
get_file_size = _es_utils.get_file_size

from config import EXCEL_MAX_ROWS_PER_SHEET, EXCEL_CHUNK_ROWS

logger = setup_logging(__name__)


class DocumentProcessor:
    
    def __init__(self, chunk_size: int = 1000, chunk_overlap: int = 200):
        self.chunk_size = chunk_size
        self.chunk_overlap = chunk_overlap
        
        self.text_splitter = RecursiveCharacterTextSplitter(
            chunk_size=chunk_size,
            chunk_overlap=chunk_overlap,
            length_function=len,
            separators=["\n\n", "\n", ". ", "! ", "? ", "; ", ": ", " ", ""],
            is_separator_regex=False
        )
        
        logger.info(f"DocumentProcessor initialized (chunk_size={chunk_size}, overlap={chunk_overlap})")
        logger.info(f"Using RecursiveCharacterTextSplitter for semantic chunking")
        self.global_chunk_counter = 0
        self.current_chunk_counter = None
        self.current_doc_id = None
    
    def _generate_unique_id(self) -> str:
        return str(uuid.uuid4())
    
    def _generate_doc_id(self, file_path: Path) -> str:
        """Generate unique document ID using SHA256 hash."""
        with open(file_path, 'rb') as f:
            file_hash = hashlib.sha256(f.read()).hexdigest()
        return f"{file_path.name}::{file_hash}"
    
    def _get_next_chunk_id(self) -> int:
        if getattr(self, 'current_chunk_counter', None) is not None:
            cid = self.current_chunk_counter
            self.current_chunk_counter += 1
            return cid

        chunk_id = self.global_chunk_counter
        self.global_chunk_counter += 1
        return chunk_id
    
    def _check_text_quality(self, text: str) -> tuple[bool, str]:
        """
        Check if text is of acceptable quality (not corrupted OCR).
        
        Detects patterns like "T he la y er" (broken words from bad PDF extraction).
        
        Returns:
            (is_good_quality, reason)
        """
        if not text or len(text.strip()) < 20:
            return False, "Text too short"
        
        # Count single-character "words" (excluding common ones like 'I', 'a')
        words = text.split()
        if not words:
            return False, "No words found"
        
        single_char_words = [w for w in words if len(w) == 1 and w.lower() not in ('i', 'a', '1', '2', '3', '4', '5', '6', '7', '8', '9', '0', '-', '/', '(', ')', '.')]
        single_char_ratio = len(single_char_words) / len(words)
        
        # If more than 15% of words are single characters, it's likely bad OCR
        if single_char_ratio > 0.15:
            return False, f"Too many single-char words ({single_char_ratio:.1%}), likely bad OCR"
        
        # Check for consecutive single-char patterns like "T h e"
        consecutive_singles = 0
        max_consecutive = 0
        for word in words:
            if len(word) == 1:
                consecutive_singles += 1
                max_consecutive = max(max_consecutive, consecutive_singles)
            else:
                consecutive_singles = 0
        
        if max_consecutive >= 4:
            return False, f"Found {max_consecutive} consecutive single-char words (broken text)"
        
        return True, "OK"
    
    
    def process_file(self, file_path: str) -> List[Dict[str, Any]]:
        file_path = Path(file_path)
        
        if not file_path.exists():
            logger.error(f"File not found: {file_path}")
            return []
        
        file_ext = file_path.suffix.lower()
        filename = file_path.name
        self.current_doc_id = self._generate_doc_id(file_path)

        logger.info(f"Processing: {filename}")
        self.current_chunk_counter = 0
        
        try:
            if file_ext == '.pdf':
                return self._process_pdf(file_path)
            elif file_ext in ['.docx', '.doc']:
                return self._process_docx(file_path)
            elif file_ext == '.txt':
                return self._process_txt(file_path)
            elif file_ext == '.csv':
                return self._process_csv(file_path)
            elif file_ext in ['.xlsx', '.xls']:
                return self._process_excel(file_path)
            elif file_ext in ['.pptx', '.ppt']:
                return self._process_pptx(file_path)
            else:
                logger.warning(f"Unsupported file type: {file_ext}")
                return []
        
        except Exception as e:
            logger.error(f"Error processing {filename}: {str(e)}")
            return []
        finally:
            self.current_doc_id = None
    
    def _process_pdf(self, file_path: Path) -> List[Dict[str, Any]]:
        documents = []

        # --- NEW: small rolling tail for cross-page semantic continuity
        previous_page_tail = ""
        TAIL_CHAR_LIMIT = 300  # safe, bounded, low risk

        try:
            pdf_document = fitz.open(file_path)
            total_pages = len(pdf_document)

            logger.info(f"PDF has {total_pages} pages")

            for page_num in range(total_pages):
                page = pdf_document[page_num]
                raw_text = page.get_text("text")
                current_text = clean_text(raw_text)

                if not current_text.strip():
                    logger.warning(f"Page {page_num + 1} is empty")
                    previous_page_tail = ""
                    continue

                # --- NEW: inject previous page tail for semantic continuity
                if previous_page_tail:
                    text_for_chunking = previous_page_tail + "\n" + current_text
                    used_cross_page_context = True
                else:
                    text_for_chunking = current_text
                    used_cross_page_context = False

                page_chunks = self.text_splitter.split_text(text_for_chunking)

                for chunk_idx, chunk in enumerate(page_chunks):
                    # Quality check - skip badly-extracted text
                    is_good, reason = self._check_text_quality(chunk)
                    if not is_good:
                        logger.warning(f"Skipping chunk from page {page_num + 1}: {reason}")
                        continue
                    
                    chunk_position = (
                        "beginning" if chunk_idx == 0 else
                        "end" if chunk_idx == len(page_chunks) - 1 else
                        "middle"
                    )

                    # --- NEW: page_range support
                    if used_cross_page_context:
                        page_range = f"{page_num}-{page_num + 1}"
                    else:
                        page_range = str(page_num + 1)

                    doc = {
                        'doc_id': self.current_doc_id,
                        'unique_id': self._generate_unique_id(),
                        'filename': file_path.name,
                        'file_path': str(file_path.absolute()),
                        'file_type': 'PDF',
                        'file_size': get_file_size(str(file_path)),
                        'page_number': page_num + 1,
                        'page_range': page_range,
                        'chunk_id': self._get_next_chunk_id(),
                        'total_pages': total_pages,
                        'total_chunks': len(page_chunks),
                        'chunk_position': chunk_position,
                        'content': chunk,
                        'chunk_text': chunk,
                        'word_count': count_words(chunk)
                    }

                    documents.append(doc)

                # --- NEW: update tail from *current* page only
                previous_page_tail = current_text[-TAIL_CHAR_LIMIT:]

            pdf_document.close()
            logger.info(
                f"Extracted {len(documents)} chunks from {total_pages} pages "
                f"(page-wise chunking with semantic cross-page continuity)"
            )

        except Exception as e:
            logger.error(f"PDF processing error: {str(e)}")

        return documents

    
    def _process_docx(self, file_path: Path) -> List[Dict[str, Any]]:
        documents = []
        
        try:
            doc = Document(file_path)
            paragraphs = [para.text for para in doc.paragraphs if para.text.strip()]
            
            if not paragraphs:
                logger.warning(f"No text found in {file_path.name}")
                return []
            
            full_text = '\n\n'.join(paragraphs)
            full_text = clean_text(full_text)
            
            logger.info(f"DOCX has {len(paragraphs)} paragraphs, {count_words(full_text)} words")
            
            chunks = self.text_splitter.split_text(full_text)
            
            for chunk_idx, chunk in enumerate(chunks):
                chunk_position = "beginning" if chunk_idx == 0 else ("end" if chunk_idx == len(chunks) - 1 else "middle")
                
                doc = {
                    'doc_id': self.current_doc_id,
                    'unique_id': self._generate_unique_id(),
                    'filename': file_path.name,
                    'file_path': str(file_path.absolute()),
                    'file_type': 'DOCX',
                    'file_size': get_file_size(str(file_path)),
                    'chunk_id': self._get_next_chunk_id(),
                    'total_chunks': len(chunks),
                    'chunk_position': chunk_position,
                    'content': chunk,
                    'chunk_text': chunk,
                    'word_count': count_words(chunk),
                    'paragraph_count': len(paragraphs)
                }
                documents.append(doc)
            
            logger.info(f"Extracted {len(documents)} chunks from DOCX")
            
        except Exception as e:
            logger.error(f"DOCX processing error: {str(e)}")
        
        return documents
    
    def _process_txt(self, file_path: Path) -> List[Dict[str, Any]]:
        documents = []
        
        try:
            encodings = ['utf-8', 'latin-1', 'cp1252', 'iso-8859-1']
            text = None
            encoding_used = None
            
            for encoding in encodings:
                try:
                    with open(file_path, 'r', encoding=encoding) as file:
                        text = file.read()
                    encoding_used = encoding
                    break
                except UnicodeDecodeError:
                    continue
            
            if text is None:
                logger.error(f"Could not decode {file_path.name} with any encoding")
                return []
            
            if not text.strip():
                logger.warning(f"Empty file: {file_path.name}")
                return []
            
            text = clean_text(text)
            line_count = text.count('\n') + 1
            
            logger.info(f"TXT has {line_count} lines, {count_words(text)} words (encoding: {encoding_used})")
            
            chunks = self.text_splitter.split_text(text)
            
            for chunk_idx, chunk in enumerate(chunks):
                chunk_position = "beginning" if chunk_idx == 0 else ("end" if chunk_idx == len(chunks) - 1 else "middle")
                
                doc = {
                    'doc_id': self.current_doc_id,
                    'unique_id': self._generate_unique_id(),
                    'filename': file_path.name,
                    'file_path': str(file_path.absolute()),
                    'file_type': 'TXT',
                    'file_size': get_file_size(str(file_path)),
                    'chunk_id': self._get_next_chunk_id(),
                    'total_chunks': len(chunks),
                    'chunk_position': chunk_position,
                    'content': chunk,
                    'chunk_text': chunk,
                    'word_count': count_words(chunk),
                    'line_count': line_count
                }
                documents.append(doc)
            
            logger.info(f"Extracted {len(documents)} chunks from TXT")
            
        except Exception as e:
            logger.error(f"TXT processing error: {str(e)}")
        
        return documents
    
    def _process_csv(self, file_path: Path) -> List[Dict[str, Any]]:
        documents = []

        try:
            df = pd.read_csv(file_path)

            if df.empty:
                logger.warning(f"Empty CSV: {file_path.name}")
                return []

            logger.info(f"CSV has {len(df)} rows, {len(df.columns)} columns")

            column_names = list(df.columns)

            # --- NEW: dataset-level summary chunk (Version-1 strength restored)
            summary_text = (
                f"Dataset: {file_path.name}\n"
                f"Columns: {', '.join(column_names)}\n"
                f"Total Rows: {len(df)}"
            )

            documents.append({
                'doc_id': self.current_doc_id,
                'unique_id': self._generate_unique_id(),
                'filename': file_path.name,
                'file_path': str(file_path.absolute()),
                'file_type': 'CSV',
                'file_size': get_file_size(str(file_path)),
                'chunk_id': self._get_next_chunk_id(),
                'total_chunks': len(df) + 1,
                'content': summary_text,
                'chunk_text': summary_text,
                'word_count': count_words(summary_text),
                'is_dataset_summary': True
            })

            column_info = f"Columns: {', '.join(column_names)}"

            for idx, row in df.iterrows():
                row_text = ' | '.join(
                    [f"{col}: {val}" for col, val in row.items() if pd.notna(val)]
                )
                row_text = clean_text(row_text)

                full_content = f"{column_info}\nRow {idx + 1}: {row_text}"

                doc = {
                    'doc_id': self.current_doc_id,
                    'unique_id': self._generate_unique_id(),
                    'filename': file_path.name,
                    'file_path': str(file_path.absolute()),
                    'file_type': 'CSV',
                    'file_size': get_file_size(str(file_path)),
                    'chunk_id': self._get_next_chunk_id(),
                    'total_chunks': len(df) + 1,
                    'content': full_content,
                    'chunk_text': full_content,
                    'word_count': count_words(full_content),
                    'row_number': int(idx) + 1,
                    'column_names': column_names,
                    'is_data_row': True
                }
                documents.append(doc)

            logger.info(
                f"Extracted {len(documents)} chunks from CSV "
                f"(1 dataset summary + {len(df)} rows)"
            )

        except Exception as e:
            logger.error(f"CSV processing error: {str(e)}")

        return documents

    
    def _process_excel(self, file_path: Path) -> List[Dict[str, Any]]:
        documents = []
        excel_file = None

        try:
            excel_file = pd.ExcelFile(file_path)
            sheet_names = excel_file.sheet_names

            logger.info(f"Excel has {len(sheet_names)} sheets: {sheet_names}")

            for sheet_name in sheet_names:
                try:
                    df = excel_file.parse(sheet_name=sheet_name, nrows=0)
                    total_rows = len(excel_file.parse(sheet_name=sheet_name))
                    
                    if total_rows == 0:
                        logger.warning(f"Empty sheet: {sheet_name}")
                        continue
                    
                    if total_rows > EXCEL_MAX_ROWS_PER_SHEET:
                        logger.warning(f"Sheet '{sheet_name}' has {total_rows} rows (limit: {EXCEL_MAX_ROWS_PER_SHEET}), processing first {EXCEL_MAX_ROWS_PER_SHEET} rows")
                        total_rows = EXCEL_MAX_ROWS_PER_SHEET
                    
                    logger.info(f"Sheet '{sheet_name}': {total_rows} rows")
                    
                    # Get column names from first row
                    df_header = excel_file.parse(sheet_name=sheet_name, nrows=0)
                    column_names = list(df_header.columns)
                    
                    for chunk_start in range(0, total_rows, EXCEL_CHUNK_ROWS):
                        chunk_end = min(chunk_start + EXCEL_CHUNK_ROWS, total_rows)
                        
                        # Always skip header row (row 0) and read from data rows
                        # skiprows=range(1, chunk_start+1) skips rows 1 to chunk_start (keeping header)
                        if chunk_start == 0:
                            df_chunk = excel_file.parse(
                                sheet_name=sheet_name,
                                nrows=chunk_end - chunk_start
                            )
                        else:
                            # For subsequent chunks, we need to re-read with header and skip processed rows
                            df_chunk = excel_file.parse(
                                sheet_name=sheet_name,
                                skiprows=range(1, chunk_start + 1),  # Skip header counts as row 0, data starts at 1
                                nrows=chunk_end - chunk_start
                            )
                            # Restore column names since skiprows may affect header detection
                            df_chunk.columns = column_names
                        
                        if df_chunk.empty:
                            continue
                        
                        column_names = list(df_chunk.columns)
                        column_info = f"Sheet: {sheet_name} | Columns: {', '.join(column_names)}"
                        
                        for idx, row in df_chunk.iterrows():
                            row_text = ' | '.join([f"{col}: {val}" for col, val in row.items() if pd.notna(val)])
                            row_text = clean_text(row_text)
                            
                            full_content = f"{column_info}\nRow {chunk_start + idx + 1}: {row_text}"
                            
                            doc = {
                                'doc_id': self.current_doc_id,
                                'unique_id': self._generate_unique_id(),
                                'filename': file_path.name,
                                'file_path': str(file_path.absolute()),
                                'file_type': 'EXCEL',
                                'file_size': get_file_size(str(file_path)),
                                'chunk_id': self._get_next_chunk_id(),
                                'total_chunks': total_rows,
                                'content': full_content,
                                'chunk_text': full_content,
                                'word_count': count_words(full_content),
                                'row_number': int(chunk_start + idx) + 1,
                                'sheet_name': sheet_name,
                                'column_names': column_names,
                                'is_data_row': True
                            }
                            documents.append(doc)
                
                except Exception as sheet_error:
                    logger.error(f"Error processing sheet '{sheet_name}': {sheet_error}")
                    continue

            logger.info(f"Extracted {len(documents)} rows from Excel (all sheets, chunked reading)")

        except Exception as e:
            logger.error(f"Excel processing error: {str(e)}")

        finally:
            try:
                if excel_file is not None:
                    excel_file.close()
                    del excel_file
            except Exception as close_err:
                logger.warning(f"Failed to close Excel file handle for {file_path.name}: {close_err}")

        return documents
    
    def _process_pptx(self, file_path: Path) -> List[Dict[str, Any]]:
        documents = []
        
        try:
            prs = Presentation(file_path)
            total_slides = len(prs.slides)
            
            logger.info(f"PPTX has {total_slides} slides")
            
            for slide_num, slide in enumerate(prs.slides, 1):
                slide_text = ""
                
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        slide_text += shape.text + "\n"
                
                if not slide_text.strip():
                    logger.warning(f"Slide {slide_num} has no text")
                    continue
                
                slide_text = clean_text(slide_text)
                chunks = self.text_splitter.split_text(slide_text)
                
                for chunk_idx, chunk in enumerate(chunks):
                    chunk_position = "beginning" if chunk_idx == 0 else ("end" if chunk_idx == len(chunks) - 1 else "middle")
                    
                    doc = {
                        'doc_id': self.current_doc_id,
                        'unique_id': self._generate_unique_id(),
                        'filename': file_path.name,
                        'file_path': str(file_path.absolute()),
                        'file_type': 'PPTX',
                        'file_size': get_file_size(str(file_path)),
                        'slide_number': slide_num,
                        'chunk_id': self._get_next_chunk_id(),
                        'total_slides': total_slides,
                        'total_chunks': len(chunks),
                        'chunk_position': chunk_position,
                        'content': chunk,
                        'chunk_text': chunk,
                        'word_count': count_words(chunk)
                    }
                    documents.append(doc)
            
            logger.info(f"Extracted {len(documents)} chunks from {total_slides} slides")
            
        except Exception as e:
            logger.error(f"PPTX processing error: {str(e)}")
        
        return documents
    
    def process_multiple_files(self, file_paths: List[str]) -> Dict[str, Any]:
        all_documents = []
        stats = {
            'total_files': len(file_paths),
            'successful': 0,
            'failed': 0,
            'total_documents': 0,
            'by_type': {}
        }
        
        logger.info(f"Processing {len(file_paths)} files")
        
        for file_path in file_paths:
            try:
                docs = self.process_file(file_path)
                
                if docs:
                    all_documents.extend(docs)
                    stats['successful'] += 1
                    stats['total_documents'] += len(docs)
                    
                    file_type = docs[0]['file_type']
                    if file_type not in stats['by_type']:
                        stats['by_type'][file_type] = 0
                    stats['by_type'][file_type] += len(docs)
                else:
                    stats['failed'] += 1
                    
            except Exception as e:
                logger.error(f"Failed to process {file_path}: {str(e)}")
                stats['failed'] += 1
        
        logger.info(f"Processing complete: {stats['successful']}/{stats['total_files']} files, {stats['total_documents']} documents")
        
        return {
            'documents': all_documents,
            'stats': stats
        }
    
    def process_uploaded_file(self, uploaded_file, save_dir: str = "data/temp") -> List[Dict[str, Any]]:
        try:
            os.makedirs(save_dir, exist_ok=True)
            
            file_path = os.path.join(save_dir, uploaded_file.name)
            with open(file_path, "wb") as f:
                f.write(uploaded_file.getbuffer())
            
            logger.info(f"Saved uploaded file: {uploaded_file.name}")
            
            documents = self.process_file(file_path)
            
            return documents
            
        except Exception as e:
            logger.error(f"Error processing uploaded file {uploaded_file.name}: {str(e)}")
            return []


if __name__ == "__main__":
    
    print("=" * 70)
    print("DOCUMENT PROCESSOR TEST")
    print("=" * 70)
    
    processor = DocumentProcessor(chunk_size=1000, chunk_overlap=200)
    
    test_files = [
        "data/temp/Perceptron.pdf",
        "data/temp/sample.docx",
        "data/temp/Data profiling.txt",
        "data/temp/action.csv",
        "data/temp/Attr-POC Approach 1.xlsx"
    ]
    
    available_files = [f for f in test_files if os.path.exists(f)]
    
    if available_files:
        print(f"\nFound {len(available_files)} test files")
        
        result = processor.process_multiple_files(available_files)
        
        print(f"\n Results:")
        print(f"  - Total files processed: {result['stats']['total_files']}")
        print(f"  - Successful: {result['stats']['successful']}")
        print(f"  - Failed: {result['stats']['failed']}")
        print(f"  - Total documents: {result['stats']['total_documents']}")
        print(f"\n By file type:")
        for file_type, count in result['stats']['by_type'].items():
            print(f"  - {file_type}: {count} documents")
        
        if result['documents']:
            print(f"\n Sample document (first chunk):")
            sample = result['documents'][0]
            for key, value in sample.items():
                if key == 'content':
                    print(f"  - {key}: {str(value)[:100]}...")
                elif key == 'csv_data':
                    print(f"  - {key}: {str(value)[:100]}...")
                else:
                    print(f"  - {key}: {value}")
    else:
        print(f"\n No test files found")
        print("Please create a 'data' folder and add sample files:")
        for f in test_files:
            print(f"  - {f}")
    
    print("\n" + "=" * 70)